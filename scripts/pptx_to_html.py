#!/usr/bin/env python3
"""
Convert a .pptx file to an HTML slideshow, and optionally create per-slide editable HTML pages.

Requirements:
 - LibreOffice installed and on PATH (for conversion to PNG)
 - python-pptx (pip install python-pptx)

Usage:
  python3 scripts/pptx_to_html.py path/to/presentation.pptx -o output.html
  python3 scripts/pptx_to_html.py path/to/presentation.pptx --per-slide --pages-dir pages --slides-dir slides

This script converts the slides into PNG files using LibreOffice and extracts slide notes using python-pptx, then writes an HTML slideshow including those images and notes. It can also create per-slide pages with overlay editable text boxes to allow static editing in-browser.
"""

import os
import sys
import argparse
import subprocess
from pptx import Presentation
from pathlib import Path
import shutil
import base64


TEMPLATE_PATH = Path(__file__).resolve().parents[1] / 'football_slideshow_template.html'
HTML_PLACEHOLDER = '<!-- Slides will be inserted by the generator or can be manually placed here. -->'


def convert_with_libreoffice(pptx_path, outdir: Path):
    p = subprocess.run([
        'soffice', '--headless', '--convert-to', 'png', '--outdir', str(outdir), str(pptx_path)
    ], capture_output=True, text=True)
    if p.returncode != 0:
        print('LibreOffice conversion failed:', p.stderr)
        raise RuntimeError('LibreOffice convert failed')
    return list(outdir.glob('*.png'))


def gather_notes(pptx_path):
    prs = Presentation(str(pptx_path))
    notes = []
    for slide in prs.slides:
        note_text = ''
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        note_text += run.text
            note_text = note_text.strip()
        notes.append(note_text)
    return notes


def build_html(slide_images, notes, out_html_path: Path, template_html: str):
    slides_html = []
    slide_images_sorted = sorted(slide_images, key=lambda p: p.name)
    for i, img_path in enumerate(slide_images_sorted):
        note = (notes[i] if i < len(notes) else '')
        rel_path = os.path.join(Path(img_path).parent.name, Path(img_path).name)
        slides_html.append(f'    <div class="slide" data-notes="{escape_attr(note)}">\n      <img class="slide-image" src="{rel_path}" alt="Slide {i+1}"/>\n    </div>')

    slides_block = '\n'.join(slides_html)
    content = template_html.replace(HTML_PLACEHOLDER, slides_block)
    out_html_path.write_text(content, encoding='utf-8')
    print('Wrote:', out_html_path)


def escape_attr(s: str) -> str:
    s = s.replace('&', '&amp;')
    s = s.replace('<', '&lt;')
    s = s.replace('>', '&gt;')
    s = s.replace('"', '&quot;')
    s = s.replace("'", '&#39;')
    return s


def build_per_slide_pages(pptx_path, slide_images, out_pages_dir: Path):
    prs = Presentation(str(pptx_path))
    slide_images_sorted = sorted(slide_images, key=lambda p: p.name)
    out_pages_dir.mkdir(parents=True, exist_ok=True)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    page_files = []
    for i, slide in enumerate(prs.slides):
        img_path = None
        if i < len(slide_images_sorted):
            img_path = slide_images_sorted[i]
        page_name = f"slide-{i+1:02d}.html"
        page_path = out_pages_dir / page_name
        # Build overlay elements for text shapes
        overlays = []
        for shape in slide.shapes:
            if not hasattr(shape, 'left'):
                continue
            left = getattr(shape, 'left', 0)
            top = getattr(shape, 'top', 0)
            width = getattr(shape, 'width', 0)
            height = getattr(shape, 'height', 0)
            left_pct = left/slide_w * 100 if slide_w else 0
            top_pct = top/slide_h * 100 if slide_h else 0
            width_pct = width/slide_w * 100 if slide_w else 0
            height_pct = height/slide_h * 100 if slide_h else 0
            style = f'position:absolute;left:{left_pct:.3f}%;top:{top_pct:.3f}%;width:{width_pct:.3f}%;height:{height_pct:.3f}%;'
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = escape_attr(shape.text or '')
                # prefer to set font-size if available from first run
                font_sz = None
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font and run.font.size:
                            font_sz = run.font.size.pt if hasattr(run.font.size, 'pt') else run.font.size
                            break
                    if font_sz:
                        break
                font_style = f'font-size:{int(font_sz)}px;' if font_sz else ''
                element = f'<div class="overlay" contenteditable="true" style="{style}{font_style}">{text}</div>'
                overlays.append(element)

        # Build the page content
        html = [
            '<!doctype html>',
            '<html lang="en">',
            '<head>',
            '  <meta charset="utf-8">',
            '  <meta name="viewport" content="width=device-width, initial-scale=1">',
            f'  <title>Slide {i+1}</title>',
            '  <style>',
            '    body{margin:0;background:#222;color:#fff;font-family:Arial,Helvetica,sans-serif}',
            '    .slide-wrap{position:relative;width:100%;height:100vh;display:flex;align-items:center;justify-content:center;}',
            '    .slide-image{max-width:100%;max-height:100%;display:block}',
            '    .overlay{position:absolute;background:transparent;color:#fff;outline:1px dashed rgba(255,255,255,0.12);padding:4px;box-sizing:border-box;white-space:pre-wrap;}',
            '    .toolbar{position:fixed;left:10px;top:10px;display:flex;gap:6px;border-radius:6px}',
            '    button{padding:8px 12px;border-radius:6px;border:1px solid #444;background:#333;color:#fff;cursor:pointer}',
            '  </style>',
            '</head>',
            '<body>',
            '  <div class="toolbar">',
            '    <button id="download">Download HTML</button>',
            '  </div>',
            '  <div class="slide-wrap">',
        ]
        if img_path is not None:
            # Read and base64-embed the image so the page is self-contained
            with open(img_path, 'rb') as f:
                data = f.read()
            b64 = base64.b64encode(data).decode('ascii')
            mime = 'image/png'
            html.append(f'    <img class="slide-image" src="data:{mime};base64,{b64}" alt="Slide {i+1}">')
        html.append('    <div class="overlay-container">')
        html.extend([f'      {o}' for o in overlays])
        html.append('    </div>')
        html.append('  </div>')
        html.append('  <script>')
        html.append('    document.getElementById("download").addEventListener("click", function(){')
        html.append('      const html = "<!doctype html>\n" + document.documentElement.outerHTML')
        html.append('      const blob = new Blob([html], {type: "text/html"})')
        html.append('      const url = URL.createObjectURL(blob)')
        html.append('      const a = document.createElement("a")')
        html.append('      a.href = url; a.download = "slide-' + f'{i+1:02d}' + '.html"; a.click()')
        html.append('    })')
        # No prev/next navigation on per-slide pages — they are independent
        html.append('  </script>')
        html.append('</body>')
        html.append('</html>')

        page_path.write_text('\n'.join(html), encoding='utf-8')
        print('Wrote per-slide page:', page_path)
        page_files.append((page_name, img_path))

    # Build index page linking to per-slide pages
    index = ['<!doctype html>', '<html lang="en">', '<head>', '  <meta charset="utf-8">', '  <meta name="viewport" content="width=device-width, initial-scale=1">', '  <title>Slides index</title>', '  <style>body{font-family:Arial,Helvetica,sans-serif;background:#111;color:#fff;padding:18px} .thumb{display:inline-block;margin:6px;border:1px solid #444;padding:6px} img{max-width:240px;height:auto;display:block} a{color:#e6e6e6;text-decoration:none}</style>', '</head>', '<body>', '<h1>Slides</h1>', '<div>']
    for name, img in page_files:
        img_rel = os.path.join(img.parent.name, img.name) if img else ''
        index.append(f'<div class="thumb"><a href="{name}"><img src="{img_rel}" alt="{name}"><div>{name}</div></a></div>')
    index.append('</div></body></html>')
    idx_path = out_pages_dir / 'index.html'
    idx_path.write_text('\n'.join(index), encoding='utf-8')
    print('Wrote pages index:', idx_path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('pptx', help='Path to the PowerPoint (.pptx) file')
    parser.add_argument('-o', '--output', default='football_slideshow.html', help='Output html file')
    parser.add_argument('--slides-dir', default='slides', help='Directory to write slide images')
    parser.add_argument('--per-slide', action='store_true', help='Create per-slide editable HTML pages in pages-dir')
    parser.add_argument('--pages-dir', default='pages', help='Directory for per-slide html pages')
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    outdir = Path(args.slides_dir).resolve()
    outdir.mkdir(parents=True, exist_ok=True)

    # Check for LibreOffice
    if not shutil.which('soffice'):
        print('LibreOffice (soffice) not found on PATH. Please install LibreOffice.')
        sys.exit(1)

    print('Converting PPTX slides to PNG images in', outdir)
    slide_images = convert_with_libreoffice(pptx_path, outdir)

    print('Reading slide notes')
    notes = gather_notes(pptx_path)

    # Load template
    if not TEMPLATE_PATH.exists():
        print('Template not found at', TEMPLATE_PATH)
        sys.exit(1)
    template_html = TEMPLATE_PATH.read_text(encoding='utf-8')

    # Build HTML
    out_html_path = Path(args.output).resolve()
    build_html(slide_images, notes, out_html_path, template_html)

    if args.per_slide:
        pages_dir = Path(args.pages_dir).resolve()
        build_per_slide_pages(pptx_path, slide_images, pages_dir)


if __name__ == '__main__':
    main()
